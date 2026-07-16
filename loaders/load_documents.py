from pathlib import Path
import json
import pandas as pd

from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    CSVLoader,
    UnstructuredHTMLLoader,
    TextLoader,
)
from langchain_core.documents import Document
from pptx import Presentation


def load_pdf(file_path):
    return PyPDFLoader(str(file_path)).load()


def load_docx(file_path):
    return Docx2txtLoader(str(file_path)).load()


def load_csv(file_path):
    return CSVLoader(str(file_path)).load()


def load_html(file_path):
    return UnstructuredHTMLLoader(str(file_path)).load()


def load_markdown(file_path):
    return TextLoader(str(file_path), encoding="utf-8").load()


def load_json(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    return [
        Document(
            page_content=json.dumps(data, indent=2),
            metadata={"source": str(file_path)},
        )
    ]


def load_excel(file_path):
    df = pd.read_excel(file_path)

    text = df.to_string(index=False)

    return [
        Document(
            page_content=text,
            metadata={"source": str(file_path)},
        )
    ]


def load_pptx(file_path):
    presentation = Presentation(file_path)

    slides_text = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(shape.text)

    text = "\n".join(slides_text)

    return [
        Document(
            page_content=text,
            metadata={"source": str(file_path)},
        )
    ]


def load_all_documents():
    documents_path = Path("documents")

    all_documents = []

    for file_path in documents_path.rglob("*"):
        if file_path.is_file():
            suffix = file_path.suffix.lower()

            try:
                if suffix == ".pdf":
                    docs = load_pdf(file_path)

                elif suffix == ".docx":
                    docs = load_docx(file_path)

                elif suffix == ".csv":
                    docs = load_csv(file_path)

                elif suffix == ".html":
                    docs = load_html(file_path)

                elif suffix == ".md":
                    docs = load_markdown(file_path)

                elif suffix == ".json":
                    docs = load_json(file_path)

                elif suffix == ".xlsx":
                    docs = load_excel(file_path)

                elif suffix == ".pptx":
                    docs = load_pptx(file_path)

                else:
                    continue

                all_documents.extend(docs)

                print(f"✓ Loaded: {file_path.name}")

            except Exception as e:
                print(f"✗ Error loading {file_path.name}: {e}")

    return all_documents


if __name__ == "__main__":
    documents = load_all_documents()

    print(f"\nTotal documents loaded: {len(documents)}")

    if documents:
        print("\nFirst document preview:\n")
        print(documents[0].page_content[:500])